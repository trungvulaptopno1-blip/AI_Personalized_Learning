# agents/content_agent/content_extractor.py
"""
Trích xuất nội dung text thô từ file upload
Hỗ trợ: PDF, DOCX, PPTX
"""

import io
import logging
from typing import Tuple, Dict, Any

try:
    import pdfplumber
    from docx import Document
    from pptx import Presentation
except ImportError as e:
    logging.error(f"Missing library for file extraction: {e}")
    raise

logger = logging.getLogger(__name__)

def extract_text_from_file(file_content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
    """
    Hàm chính: Trích xuất text và metadata từ file
    
    Returns:
        (raw_text: str, metadata: dict)
    """
    file_type = _get_file_type(filename).lower()
    metadata = {
        "filename": filename,
        "file_type": file_type,
        "size_bytes": len(file_content),
    }

    try:
        if file_type == "pdf":
            text, extra = _extract_pdf(file_content)
        elif file_type in ("docx", "doc"):
            text, extra = _extract_docx(file_content)
        elif file_type in ("pptx", "ppt"):
            text, extra = _extract_pptx(file_content)
        else:
            raise ValueError(f"Định dạng file không được hỗ trợ: {file_type}")

        metadata.update(extra)
        logger.info(f"Extracted {len(text)} chars from {filename} ({file_type})")

        return text.strip(), metadata

    except Exception as e:
        logger.error(f"Extraction failed for {filename}: {str(e)}", exc_info=True)
        raise


def _get_file_type(filename: str) -> str:
    """Lấy phần mở rộng file"""
    if '.' not in filename:
        return "unknown"
    return filename.rsplit('.', 1)[-1].lower()


def _extract_pdf(content: bytes) -> Tuple[str, Dict[str, Any]]:
    meta = {}
    parts = []

    with pdfplumber.open(io.BytesIO(content)) as pdf:
        meta["page_count"] = len(pdf.pages)
        meta["title"] = pdf.metadata.get("Title", "").strip() if pdf.metadata else ""

        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)

    return "\n\n".join(parts), meta


def _extract_docx(content: bytes) -> Tuple[str, Dict[str, Any]]:
    meta = {}
    parts = []

    doc = Document(io.BytesIO(content))
    meta["paragraph_count"] = len(doc.paragraphs)

    # Tìm tiêu đề (nếu có)
    for para in doc.paragraphs:
        if para.style.name.lower().startswith("heading"):
            meta["title"] = para.text.strip()
            break

    # Thu thập nội dung
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Thêm nội dung từ bảng (nếu có)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return "\n\n".join(parts), meta


def _extract_pptx(content: bytes) -> Tuple[str, Dict[str, Any]]:
    meta = {}
    parts = []

    prs = Presentation(io.BytesIO(content))
    meta["slide_count"] = len(prs.slides)

    if prs.slides and prs.slides[0].shapes.title:
        meta["title"] = prs.slides[0].shapes.title.text.strip()

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())

        if slide_parts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_parts))

    return "\n\n".join(parts), meta